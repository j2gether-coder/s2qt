# -*- coding: utf-8 -*-
"""
office_export.py

JSON 입력으로 DOCX / PPTX 생성
- DOCX: python-docx
- PPTX: python-pptx + 템플릿 기반

PPTX 템플릿 구조(고정):
  슬라이드 1: 제목 슬라이드
    - cover_title
    - cover_scripture
    - cover_hymn

  슬라이드 2: 본문 슬라이드 (복제용)
    - section_title
    - section_subtitle
    - section_body

사용 예:
  python office_export.py --input sample_docx.json
  python office_export.py --input sample_pptx.json
"""

import argparse
import copy
import json
import os
import sys
from typing import Any, Dict, List

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree


# --------------------------------------------------
# 경로
# --------------------------------------------------
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTX_TEMPLATE_PATH = os.path.join(BASE_DIR, "var", "template", "qt_basic_template.pptx")


# --------------------------------------------------
# 공통 유틸
# --------------------------------------------------
def load_json(path: str) -> Dict[str, Any]:
    with open(path, "r", encoding="utf-8-sig") as f:
        return json.load(f)


def ensure_parent_dir(path: str) -> None:
    parent = os.path.dirname(os.path.abspath(path))
    if parent and not os.path.exists(parent):
        os.makedirs(parent, exist_ok=True)


def safe_str(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def safe_list_of_str(value: Any) -> List[str]:
    if not isinstance(value, list):
        return []
    out: List[str] = []
    for item in value:
        s = safe_str(item)
        if s:
            out.append(s)
    return out


def parse_payload(payload: Dict[str, Any]) -> Dict[str, Any]:
    if not isinstance(payload, dict):
        raise ValueError("JSON 루트가 객체가 아닙니다.")

    fmt = safe_str(payload.get("format")).lower()
    if fmt not in ("docx", "pptx"):
        raise ValueError("format 값은 docx 또는 pptx 이어야 합니다.")

    output_path = safe_str(payload.get("output_path"))
    if not output_path:
        raise ValueError("output_path 값이 없습니다.")

    document = payload.get("document")
    if not isinstance(document, dict):
        raise ValueError("document 객체가 없습니다.")

    return {
        "format": fmt,
        "output_path": output_path,
        "document": document,
    }


def join_paragraphs(paragraphs: List[str]) -> str:
    return "\n\n".join([p for p in paragraphs if safe_str(p)])


def build_numbered_lines(items: List[str]) -> str:
    lines: List[str] = []
    for idx, item in enumerate(items, start=1):
        s = safe_str(item)
        if s:
            lines.append(f"{idx}. {s}")
    return "\n\n".join(lines)


def build_qt_document(data: Dict[str, Any]) -> Dict[str, Any]:
    """
    temp.json 구조를 내부 문서 모델로 변환
    """
    title = safe_str(data.get("title"))
    bible_text = safe_str(data.get("bible_text"))
    hymn = safe_str(data.get("hymn"))

    summary = safe_list_of_str(data.get("summary"))
    reflection_items = safe_list_of_str(data.get("reflection_items"))
    prayer_title = safe_str(data.get("prayer_title"))
    prayer_paragraphs = safe_list_of_str(data.get("prayer_paragraphs"))
    footer_text = safe_str(data.get("footer_text"))

    messages_raw = data.get("messages", [])
    messages: List[Dict[str, str]] = []
    if isinstance(messages_raw, list):
        for idx, msg in enumerate(messages_raw, start=1):
            if not isinstance(msg, dict):
                continue
            msg_title = safe_str(msg.get("title")) or f"메시지 {idx}"
            paragraphs = safe_list_of_str(msg.get("paragraphs"))
            messages.append(
                {
                    "title": msg_title,
                    "body_text": join_paragraphs(paragraphs),
                }
            )

    return {
        "title": title,
        "bible_text": bible_text,
        "hymn": hymn,
        "summary_text": join_paragraphs(summary),
        "messages": messages,
        "reflection_text": build_numbered_lines(reflection_items),
        "prayer_title": prayer_title if prayer_title else "오늘의 기도",
        "prayer_text": join_paragraphs(prayer_paragraphs),
        "footer_text": footer_text,
    }


# --------------------------------------------------
# DOCX
# --------------------------------------------------
def ensure_docx_styles(doc: Document) -> None:
    styles = doc.styles

    normal = styles["Normal"]
    normal.font.name = "맑은 고딕"
    normal.font.size = Pt(11)

    title = styles["Title"]
    title.font.name = "맑은 고딕"
    title.font.size = Pt(20)
    title.font.bold = True

    h1 = styles["Heading 1"]
    h1.font.name = "맑은 고딕"
    h1.font.size = Pt(14)
    h1.font.bold = True

    h2 = styles["Heading 2"]
    h2.font.name = "맑은 고딕"
    h2.font.size = Pt(12)
    h2.font.bold = True


def add_docx_styled_paragraph(
    doc: Document,
    text: str,
    style: str = "Normal",
    align=None,
    space_after: int = None,
) -> None:
    p = doc.add_paragraph(text, style=style)
    if align is not None:
        p.alignment = align
    if space_after is not None:
        p.paragraph_format.space_after = Pt(space_after)


def add_docx_blank(doc: Document, count: int = 1) -> None:
    for _ in range(count):
        doc.add_paragraph("")


def export_docx(output_path: str, data: Dict[str, Any]) -> None:
    qt = build_qt_document(data)

    doc = Document()
    ensure_docx_styles(doc)

    section = doc.sections[0]
    section.top_margin = Pt(56)
    section.bottom_margin = Pt(56)
    section.left_margin = Pt(56)
    section.right_margin = Pt(56)

    if qt["title"]:
        add_docx_styled_paragraph(
            doc,
            qt["title"],
            style="Title",
            align=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=10,
        )

    if qt["bible_text"]:
        add_docx_styled_paragraph(doc, f"본문: {qt['bible_text']}", style="Normal", space_after=2)
    if qt["hymn"]:
        add_docx_styled_paragraph(doc, f"찬송: {qt['hymn']}", style="Normal", space_after=2)

    if qt["bible_text"] or qt["hymn"]:
        add_docx_blank(doc, 1)

    if qt["summary_text"]:
        add_docx_styled_paragraph(doc, "말씀의 창", style="Heading 1", space_after=6)
        for para in qt["summary_text"].split("\n\n"):
            if safe_str(para):
                add_docx_styled_paragraph(doc, para, style="Normal", space_after=4)
        add_docx_blank(doc, 1)

    if qt["messages"]:
        add_docx_styled_paragraph(doc, "오늘의 메시지", style="Heading 1", space_after=6)
        for idx, msg in enumerate(qt["messages"], start=1):
            add_docx_styled_paragraph(
                doc,
                f"{idx}. {msg['title']}",
                style="Heading 2",
                space_after=4,
            )
            for para in msg["body_text"].split("\n\n"):
                if safe_str(para):
                    add_docx_styled_paragraph(doc, para, style="Normal", space_after=3)
            if idx < len(qt["messages"]):
                add_docx_blank(doc, 1)
        add_docx_blank(doc, 1)

    if qt["reflection_text"]:
        add_docx_styled_paragraph(doc, "삶의 적용", style="Heading 1", space_after=6)
        for para in qt["reflection_text"].split("\n\n"):
            if safe_str(para):
                add_docx_styled_paragraph(doc, para, style="Normal", space_after=3)
        add_docx_blank(doc, 1)

    if qt["prayer_text"]:
        add_docx_styled_paragraph(doc, qt["prayer_title"], style="Heading 1", space_after=6)
        for para in qt["prayer_text"].split("\n\n"):
            if safe_str(para):
                add_docx_styled_paragraph(doc, para, style="Normal", space_after=4)
        add_docx_blank(doc, 1)

    if qt["footer_text"]:
        add_docx_styled_paragraph(
            doc,
            qt["footer_text"],
            style="Normal",
            align=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=0,
        )

    ensure_parent_dir(output_path)
    doc.save(output_path)


# --------------------------------------------------
# PPTX 템플릿 주입
# --------------------------------------------------
def find_shape_by_name(slide, shape_name: str):
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    return None


def replace_text_keep_style(slide, shape_name: str, new_text: str) -> None:
    """
    shape 이름으로 텍스트를 완전히 교체
    - 기존 첫 문단의 pPr(정렬/문단속성) 복사
    - 기존 첫 run의 rPr(글자속성) 복사
    - 문단 구분은 \\n\\n 기준
    """
    shape = find_shape_by_name(slide, shape_name)
    if shape is None:
        raise ValueError(f"shape '{shape_name}' 를 찾을 수 없습니다.")

    if not shape.has_text_frame:
        raise ValueError(f"shape '{shape_name}' 는 text frame이 아닙니다.")

    tf = shape.text_frame
    txBody = tf._txBody
    old_paragraphs = list(txBody.findall(qn("a:p")))

    src_pPr = None
    src_rPr = None

    if old_paragraphs:
        first_p = old_paragraphs[0]

        pPr = first_p.find(qn("a:pPr"))
        if pPr is not None:
            src_pPr = copy.deepcopy(pPr)

        first_r = first_p.find(qn("a:r"))
        if first_r is not None:
            rPr = first_r.find(qn("a:rPr"))
            if rPr is not None:
                src_rPr = copy.deepcopy(rPr)

    # 기존 문단 전체 제거
    for p in old_paragraphs:
        txBody.remove(p)

    paragraphs_text = new_text.split("\n\n") if new_text else [""]

    for para_text in paragraphs_text:
        p = etree.SubElement(txBody, qn("a:p"))

        if src_pPr is not None:
            p.append(copy.deepcopy(src_pPr))

        r = etree.SubElement(p, qn("a:r"))

        if src_rPr is not None:
            r.append(copy.deepcopy(src_rPr))

        t = etree.SubElement(r, qn("a:t"))
        t.text = para_text.strip()

    # 최소 1문단 보장
    if len(txBody.findall(qn("a:p"))) == 0:
        etree.SubElement(txBody, qn("a:p"))


def duplicate_slide(prs: Presentation, slide_index: int):
    """
    기존 슬라이드를 복제해서 프레젠테이션 끝에 추가
    - source slide와 동일한 layout 사용
    """
    source = prs.slides[slide_index]
    source_layout = source.slide_layout

    new_slide = prs.slides.add_slide(source_layout)

    for shape in list(new_slide.shapes):
        el = shape.element
        el.getparent().remove(el)

    for shape in source.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


def delete_slide(prs: Presentation, slide_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id_list.remove(slides[slide_index])


def render_pptx_from_template(output_path: str, data: Dict[str, Any], template_path: str = PPTX_TEMPLATE_PATH) -> None:
    qt = build_qt_document(data)

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"PPTX 템플릿 파일이 없습니다: {template_path}")

    prs = Presentation(template_path)

    if len(prs.slides) < 2:
        raise ValueError("PPTX 템플릿은 최소 2장(제목 슬라이드 + 본문 슬라이드)이어야 합니다.")

    cover_slide = prs.slides[0]
    body_template_index = 1

    # 1) 제목 슬라이드 채우기
    replace_text_keep_style(cover_slide, "cover_title", qt["title"])
    replace_text_keep_style(cover_slide, "cover_scripture", f"본문: {qt['bible_text']}")
    replace_text_keep_style(cover_slide, "cover_hymn", f"찬송: {qt['hymn']}")

    # 2) 본문 섹션 목록 구성
    sections: List[Dict[str, str]] = []

    if qt["summary_text"]:
        sections.append(
            {
                "section_title": "말씀의 창",
                "section_subtitle": "",
                "section_body": qt["summary_text"],
            }
        )

    for idx, msg in enumerate(qt["messages"], start=1):
        sections.append(
            {
                "section_title": "오늘의 메시지",
                "section_subtitle": f"{idx}. {msg['title']}",
                "section_body": msg["body_text"],
            }
        )

    if qt["reflection_text"]:
        sections.append(
            {
                "section_title": "삶의 적용",
                "section_subtitle": "",
                "section_body": qt["reflection_text"],
            }
        )

    if qt["prayer_text"]:
        sections.append(
            {
                "section_title": qt["prayer_title"],
                "section_subtitle": "",
                "section_body": qt["prayer_text"],
            }
        )

    # 본문이 없으면 cover만 남기고 body 제거
    if not sections:
        while len(prs.slides) > 1:
            delete_slide(prs, 1)
        ensure_parent_dir(output_path)
        prs.save(output_path)
        return

    # 3) section 수에 맞게 본문 슬라이드 복제
    additional_needed = len(sections) - 1
    for _ in range(additional_needed):
        duplicate_slide(prs, body_template_index)

    # 4) 각 본문 슬라이드 채우기
    for idx, section in enumerate(sections, start=1):
        slide = prs.slides[idx]
        replace_text_keep_style(slide, "section_title", section["section_title"])
        replace_text_keep_style(slide, "section_subtitle", section["section_subtitle"])
        replace_text_keep_style(slide, "section_body", section["section_body"])

    # 5) 남는 슬라이드 제거
    expected_total = 1 + len(sections)
    while len(prs.slides) > expected_total:
        delete_slide(prs, len(prs.slides) - 1)

    ensure_parent_dir(output_path)
    prs.save(output_path)


def export_pptx(output_path: str, data: Dict[str, Any]) -> None:
    render_pptx_from_template(output_path, data, PPTX_TEMPLATE_PATH)


# --------------------------------------------------
# Main
# --------------------------------------------------
def main() -> int:
    parser = argparse.ArgumentParser(description="JSON 기준 DOCX/PPTX 생성기")
    parser.add_argument("--input", required=True, help="입력 JSON 파일 경로")
    args = parser.parse_args()

    try:
        payload = load_json(args.input)
        parsed = parse_payload(payload)

        fmt = parsed["format"]
        output_path = parsed["output_path"]
        document = parsed["document"]

        if fmt == "docx":
            export_docx(output_path, document)
        else:
            export_pptx(output_path, document)

        print(f"[OK] generated: {output_path}")
        return 0

    except Exception as e:
        print(f"[ERROR] {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())