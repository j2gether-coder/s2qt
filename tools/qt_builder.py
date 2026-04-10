"""
S2QT PPTX 템플릿 빌더 + 자동 생성기
────────────────────────────────────────────────────────────────
이 파일 하나로 두 가지 작업을 수행합니다.

  1. build_template()  → QT_묵상_템플릿_v2.pptx  생성
  2. generate_qt_pptx()→ JSON 데이터를 주입해 결과물 저장

Shape 이름 (python-pptx 에서 shape.name 으로 접근):
  Slide 1 (표지):
    cover_title          ← str  설교/묵상 제목
    cover_scripture      ← str  "📖 본문: ..."
    cover_hymn           ← str  "🎵 찬송: ..."

  Slide 2 (말씀의 창):
    scripture_body       ← str  본문 요약 (\n\n 으로 문단 구분)

  Slide 3 (오늘의 메시지 1):
    msg1_subtitle        ← str  "1. 제목"
    msg1_body            ← str  본문 문단

  Slide 4 (오늘의 메시지 2):
    msg2_subtitle        ← str
    msg2_body            ← str

  Slide 5 (삶의 적용):
    application_item_1   ← str
    application_item_2   ← str
    application_item_3   ← str

  Slide 6 (오늘의 기도):
    prayer_body          ← str
"""

from __future__ import annotations
import copy, json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# ─────────────────────────────────────────────────────────────
# 색상 팔레트
# ─────────────────────────────────────────────────────────────
class C:
    BG          = RGBColor(0xFA, 0xF7, 0xF2)   # 따뜻한 아이보리
    HEADER_BAR  = RGBColor(0x3D, 0x5C, 0x4A)   # 포레스트 그린
    ACCENT      = RGBColor(0x8B, 0x6F, 0x47)   # 브라운 골드
    ACCENT_LIGHT= RGBColor(0xED, 0xE0, 0xCC)   # 연한 크림
    TITLE_TEXT  = RGBColor(0x2C, 0x3A, 0x2E)   # 다크 그린-차콜
    BODY_TEXT   = RGBColor(0x3D, 0x3D, 0x3D)   # 본문
    SUB_TEXT    = RGBColor(0x7A, 0x7A, 0x7A)   # 보조
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    DIVIDER     = RGBColor(0xD4, 0xC5, 0xA9)
    PRAYER_BG   = RGBColor(0xF4, 0xEF, 0xE6)


FONT = "Malgun Gothic"
W_IN  = Inches(10)
H_IN  = Inches(5.625)


# ─────────────────────────────────────────────────────────────
# XML 헬퍼
# ─────────────────────────────────────────────────────────────
def _rgb_str(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _set_fill(sp, color: RGBColor | None):
    """sp(spPr) 에 단색 채우기 또는 noFill 설정"""
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    # 기존 fill 제거
    for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"), qn("a:pattFill")):
        for el in spPr.findall(tag):
            spPr.remove(el)
    if color is None:
        etree.SubElement(spPr, qn("a:noFill"))
    else:
        sf = etree.SubElement(spPr, qn("a:solidFill"))
        srgb = etree.SubElement(sf, qn("a:srgbClr"))
        srgb.set("val", _rgb_str(color))


def _set_line(sp, color: RGBColor | None, width_pt: float = 0.75):
    spPr = sp.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    if color is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
    else:
        ln = etree.SubElement(spPr, qn("a:ln"))
        ln.set("w", str(int(width_pt * 12700)))
        sf = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(sf, qn("a:srgbClr"))
        srgb.set("val", _rgb_str(color))


# ─────────────────────────────────────────────────────────────
# shape 추가 헬퍼 (python-pptx 래퍼)
# ─────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill: RGBColor | None, line: RGBColor | None = None,
             line_width: float = 0.75, name: str = ""):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if name:
        shape.name = name
    sp = shape._element
    _set_fill(sp, fill)
    _set_line(sp, line, line_width)
    return shape


def add_oval(slide, x, y, w, h, fill: RGBColor, name: str = ""):
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE_TYPE.OVAL
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if name:
        shape.name = name
    sp = shape._element
    _set_fill(sp, fill)
    _set_line(sp, None)
    return shape


def add_textbox(slide, x, y, w, h,
                text: str, font_size: float, color: RGBColor,
                bold: bool = False, align: str = "left",
                name: str = "", italic: bool = False,
                line_spacing: float | None = None):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if name:
        txBox.name = name
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # 기존 단락 제거 후 재작성
    paragraphs = text.split("\n\n")
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = {
            "left":   PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right":  PP_ALIGN.RIGHT,
        }.get(align, PP_ALIGN.LEFT)

        if line_spacing:
            p.line_spacing = Pt(line_spacing)

        run = p.add_run()
        run.text = para_text.strip()
        run.font.name      = FONT
        run.font.size      = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold      = bold
        run.font.italic    = italic

    return txBox


# ─────────────────────────────────────────────────────────────
# 공통 레이아웃 요소
# ─────────────────────────────────────────────────────────────
TOTAL_SLIDES = 6


def _add_header_bar(slide, label: str):
    add_rect(slide, 0, 0, 10, 0.65, C.HEADER_BAR, C.HEADER_BAR)
    add_rect(slide, 9.55, 0, 0.45, 0.65, C.ACCENT, C.ACCENT)
    add_textbox(slide, 0.45, 0.02, 9.0, 0.6,
                label, 10.5, C.WHITE, align="left")


def _add_section_title(slide, label: str, y: float, name: str = ""):
    add_rect(slide, 0.45, y, 0.06, 0.38, C.ACCENT, C.ACCENT)
    add_textbox(slide, 0.6, y, 8.9, 0.38,
                label, 15, C.TITLE_TEXT, bold=True, align="left", name=name)


def _add_footer(slide, page_num: int):
    add_rect(slide, 0.45, 5.2, 9.1, 0.01, C.DIVIDER, C.DIVIDER)
    add_textbox(slide, 0, 5.22, 10, 0.3,
                f"{page_num} / {TOTAL_SLIDES}", 9, C.SUB_TEXT, align="center")


def _add_content_box(slide, x, y, w, h, fill=None, line=None):
    """흰색 카드 박스 (테두리 있음)"""
    add_rect(slide, x, y, w, h,
             fill or C.WHITE,
             line or C.DIVIDER, 0.8)


# ─────────────────────────────────────────────────────────────
# 슬라이드 빌더
# ─────────────────────────────────────────────────────────────
def _build_slide1_cover(prs):
    """표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title  # suppress default

    # 배경
    add_rect(slide, 0, 0, 10, 5.625, C.BG, None)
    # 헤더 바 (두꺼운)
    add_rect(slide, 0, 0, 10, 1.1, C.HEADER_BAR, C.HEADER_BAR)
    add_rect(slide, 0, 0, 0.5, 1.1, C.ACCENT, C.ACCENT)
    add_textbox(slide, 0.65, 0.05, 9.0, 1.0, "QT  묵상", 13, C.WHITE, align="left")

    # ★ cover_title
    add_textbox(slide, 0.9, 1.3, 8.2, 1.35,
                "설교 / 묵상 제목을 입력하세요",
                28, C.TITLE_TEXT, bold=True, align="center", name="cover_title")

    # 구분선
    add_rect(slide, 3.5, 2.75, 3.0, 0.025, C.ACCENT, C.ACCENT)

    # ★ cover_scripture
    add_textbox(slide, 0.9, 2.95, 8.2, 0.45,
                "📖  본문 말씀: 여기에 입력하세요",
                13, C.ACCENT, bold=True, align="center", name="cover_scripture")

    # ★ cover_hymn
    add_textbox(slide, 0.9, 3.42, 8.2, 0.45,
                "🎵  찬송: 여기에 입력하세요",
                12, C.SUB_TEXT, align="center", name="cover_hymn")

    # 하단 바
    add_rect(slide, 0, 5.22, 10, 0.4, C.ACCENT_LIGHT, None)
    add_textbox(slide, 0, 5.25, 10, 0.3,
                "본 문서는 S2QT 자동 생성 템플릿입니다",
                9, C.SUB_TEXT, align="center")


def _build_slide2_scripture(prs):
    """말씀의 창"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, C.BG, None)
    _add_header_bar(slide, "말씀의 창  ·  본문 요약")
    _add_section_title(slide, "말씀의 창 (Scripture)", 0.85, "scripture_section_title")
    _add_content_box(slide, 0.45, 1.35, 9.1, 3.5)

    # ★ scripture_body
    add_textbox(slide, 0.62, 1.45, 8.76, 3.3,
                "본문 요약 첫째 문단을 여기에 입력하세요.\n\n"
                "둘째 문단을 여기에 입력합니다.\n\n"
                "셋째 문단을 여기에 입력합니다.",
                12.5, C.BODY_TEXT, name="scripture_body", line_spacing=20)

    _add_footer(slide, 2)


def _build_slide_message(prs, page_num: int, msg_idx: int):
    """오늘의 메시지 (공통)"""
    pfx = f"msg{msg_idx}"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, C.BG, None)
    _add_header_bar(slide, "오늘의 메시지")
    _add_section_title(slide, "오늘의 메시지", 0.85, f"{pfx}_section_title")

    # 메시지 제목 강조 박스
    add_rect(slide, 0.45, 1.35, 9.1, 0.52, C.ACCENT_LIGHT, C.DIVIDER, 0.6)
    # ★ msg1_subtitle / msg2_subtitle
    add_textbox(slide, 0.55, 1.37, 8.9, 0.48,
                f"{msg_idx}.  메시지 제목을 입력하세요",
                14.5, C.TITLE_TEXT, bold=True, align="left",
                name=f"{pfx}_subtitle")

    # 본문 박스
    _add_content_box(slide, 0.45, 1.97, 9.1, 2.88)
    # ★ msg1_body / msg2_body
    add_textbox(slide, 0.62, 2.07, 8.76, 2.7,
                "첫 번째 문단을 입력합니다.\n\n"
                "두 번째 문단을 입력합니다.\n\n"
                "세 번째 문단을 입력합니다.",
                12.5, C.BODY_TEXT, name=f"{pfx}_body", line_spacing=20)

    _add_footer(slide, page_num)


def _build_slide5_application(prs):
    """삶의 적용"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, C.BG, None)
    _add_header_bar(slide, "삶의 적용  ·  Application")
    _add_section_title(slide, "삶의 적용 (Application)", 0.85, "application_section_title")

    placeholders = [
        "첫 번째 적용 사항을 입력하세요.",
        "두 번째 적용 사항을 입력하세요.",
        "세 번째 적용 사항을 입력하세요.",
    ]
    for i, text in enumerate(placeholders):
        y = 1.42 + i * 1.08
        add_oval(slide, 0.45, y + 0.05, 0.44, 0.44, C.HEADER_BAR)
        add_textbox(slide, 0.45, y + 0.05, 0.44, 0.44,
                    str(i + 1), 12, C.WHITE, bold=True, align="center")
        _add_content_box(slide, 1.05, y, 8.5, 0.82)
        # ★ application_item_1 / _2 / _3
        add_textbox(slide, 1.15, y, 8.3, 0.82,
                    text, 12.5, C.BODY_TEXT,
                    name=f"application_item_{i+1}", line_spacing=19)

    _add_footer(slide, 5)


def _build_slide6_prayer(prs):
    """오늘의 기도"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, C.BG, None)
    _add_header_bar(slide, "오늘의 기도  ·  Prayer")
    _add_section_title(slide, "오늘의 기도 (Prayer)", 0.85, "prayer_section_title")
    add_rect(slide, 0.45, 1.35, 9.1, 3.5, C.PRAYER_BG, C.ACCENT, 0.8)
    add_textbox(slide, 0.45, 1.35, 9.1, 0.55, "🙏", 18, C.BODY_TEXT, align="center")
    # ★ prayer_body
    add_textbox(slide, 0.62, 1.9, 8.76, 2.75,
                "기도 첫 번째 문단을 입력하세요.\n\n"
                "두 번째 문단을 입력하세요.\n\n"
                "예수님의 이름으로 기도드립니다. 아멘.",
                12.5, C.BODY_TEXT, name="prayer_body",
                italic=False, line_spacing=21)
    _add_footer(slide, 6)


# ─────────────────────────────────────────────────────────────
# 템플릿 생성
# ─────────────────────────────────────────────────────────────
def build_template(output_path: str = "QT_묵상_템플릿_v2.pptx"):
    prs = Presentation()
    prs.slide_width  = W_IN
    prs.slide_height = H_IN

    _build_slide1_cover(prs)
    _build_slide2_scripture(prs)
    _build_slide_message(prs, 3, 1)
    _build_slide_message(prs, 4, 2)
    _build_slide5_application(prs)
    _build_slide6_prayer(prs)

    prs.save(output_path)
    print(f"✅ 템플릿 저장 완료 → {output_path}")
    return output_path


# ─────────────────────────────────────────────────────────────
# JSON 주입 (텍스트 교체)
# ─────────────────────────────────────────────────────────────
def _replace_text(slide, shape_name: str, new_text: str):
    """named shape의 텍스트를 교체 (기존 서식 유지)"""
    shape = next((s for s in slide.shapes if s.name == shape_name), None)
    if shape is None:
        print(f"  ⚠️  shape '{shape_name}' 없음")
        return

    tf = shape.text_frame
    paragraphs_text = new_text.split("\n\n")

    # 첫 단락 서식 참조
    src_para = tf.paragraphs[0]
    src_run  = src_para.runs[0] if src_para.runs else None

    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    for para_text in paragraphs_text:
        new_p = copy.deepcopy(src_para._p)
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)

        r_elem = etree.SubElement(new_p, qn("a:r"))
        if src_run is not None:
            src_rpr = src_run._r.find(qn("a:rPr"))
            if src_rpr is not None:
                r_elem.insert(0, copy.deepcopy(src_rpr))

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = para_text.strip()
        txBody.append(new_p)

    preview = new_text[:35].replace("\n", " ")
    print(f"  ✅  [{shape_name}] ← {preview}{'...' if len(new_text) > 35 else ''}")


def generate_qt_pptx(data: dict,
                     template_path: str = "QT_묵상_템플릿_v2.pptx",
                     output_path:   str = "QT_출력_결과.pptx"):
    """
    JSON 데이터를 템플릿에 주입하여 완성된 PPTX 저장.

    data 스키마:
    {
      "title":      str,
      "scripture":  str,
      "hymns":      [str, ...],
      "summary":    str,         # \\n\\n 으로 문단 구분
      "messages": [
        { "title": str, "content": str },  # \\n\\n 으로 문단 구분
        { "title": str, "content": str }
      ],
      "reflection": [str, str, str],
      "prayer":     str          # \\n\\n 으로 문단 구분
    }
    """
    prs    = Presentation(template_path)
    slides = prs.slides

    print("📄 슬라이드 1 ─ 표지")
    s = slides[0]
    _replace_text(s, "cover_title",     data.get("title", ""))
    _replace_text(s, "cover_scripture", f"📖  본문 말씀: {data.get('scripture', '')}")
    _replace_text(s, "cover_hymn",      f"🎵  찬송: {' / '.join(data.get('hymns', []))}")

    print("\n📄 슬라이드 2 ─ 말씀의 창")
    _replace_text(slides[1], "scripture_body", data.get("summary", ""))

    msgs = data.get("messages", [{}, {}])
    for slide_idx, msg_idx in [(2, 1), (3, 2)]:
        m = msgs[msg_idx - 1] if len(msgs) >= msg_idx else {}
        print(f"\n📄 슬라이드 {slide_idx + 1} ─ 메시지 {msg_idx}")
        _replace_text(slides[slide_idx], f"msg{msg_idx}_subtitle",
                      f"{msg_idx}.  {m.get('title', '')}")
        _replace_text(slides[slide_idx], f"msg{msg_idx}_body",
                      m.get("content", ""))

    print("\n📄 슬라이드 5 ─ 삶의 적용")
    for i, text in enumerate(data.get("reflection", ["", "", ""])[:3], 1):
        _replace_text(slides[4], f"application_item_{i}", f"  {text}")

    print("\n📄 슬라이드 6 ─ 오늘의 기도")
    _replace_text(slides[5], "prayer_body", data.get("prayer", ""))

    prs.save(output_path)
    print(f"\n🎉 결과 저장 완료 → {output_path}")


# ─────────────────────────────────────────────────────────────
# 샘플 실행
# ─────────────────────────────────────────────────────────────
SAMPLE_DATA = {
    "title": "말세를 살아가는 성도의 두 가지 향기",
    "scripture": "데살로니가전서 4:9-12",
    "hymns": ["218장 (네 맘과 정성을 다하여서)"],
    "summary": (
        "사도 바울은 데살로니가 성도들에게 형제 사랑을 더욱 힘쓰라고 권면합니다.\n\n"
        "하나님의 뜻은 성도들의 거룩함에 있으며, 음란을 버리고 각자 거룩함과 존귀함으로 "
        "살아갈 것을 강조합니다.\n\n"
        "바울은 조용히 자기 일을 하며 손으로 일하기를 힘쓰라고 명령하여, 일상의 삶 속에서도 "
        "그리스도인의 향기를 드러낼 것을 촉구합니다."
    ),
    "messages": [
        {
            "title": "사랑의 용광로를 품는 공동체",
            "content": (
                "말세를 살아가는 성도의 첫 번째 자세는 형제 사랑의 열정을 회복하는 것입니다. "
                "데살로니가 교회는 이미 사랑의 수고로 소문난 곳이었지만, 바울은 더욱 그렇게 행하라고 권면합니다.\n\n"
                "성도의 사랑은 감정이 아닌 하나님의 가르침에 대한 순종입니다. "
                "내가 정한 범위를 넘어 모든 이들을 향해 마음을 여는 것이 이기주의를 이기는 성도의 모습입니다.\n\n"
                "오늘 내가 먼저 이해하고 다가가야 할 사람은 누구입니까? "
                "선입견을 내려놓고 사랑으로 다가가십시오."
            ),
        },
        {
            "title": "일상의 소명에 최선을 다하는 삶",
            "content": (
                "두 번째 자세는 소명 의식(Vocation)으로 일상에 충실한 것입니다. "
                "당시 일부 성도들은 재림을 핑계로 생업을 등한시하며 질서를 어지럽혔습니다.\n\n"
                "바울은 조용히 자기 일을 하고 너희 손으로 일하기를 힘쓰라고 명합니다. "
                "직업은 생계 수단을 넘어 하나님이 부르신 거룩한 자리입니다.\n\n"
                "코람 데오(Coram Deo)의 정신으로, 하나님의 눈앞에서 오늘 맡겨진 일에 최선을 다하는 것이 "
                "세상 속에서 그리스도의 향기를 드러내는 길입니다."
            ),
        },
    ],
    "reflection": [
        "오늘 만나는 사람들에게 따뜻한 인사를 먼저 건네고 그들의 필요를 묻겠습니다.",
        "나에게 맡겨진 일상 업무를 주님께 하듯 성실하게 처리하겠습니다.",
        "말씀과 기도, 교회 안의 사랑과 세상 속의 성실함의 균형을 회복하겠습니다.",
    ],
    "prayer": (
        "하나님 아버지, 마지막 때를 살아가는 저희가 뜨거운 사랑으로 서로를 배려하게 하옵소서.\n\n"
        "각자에게 주어진 일터와 가정에서 소명 의식을 가지고 살아가게 하시며, "
        "흑백 논리와 이기주의를 넘어 균형 잡힌 신앙으로 세상의 빛과 소금이 되게 인도하여 주시옵소서.\n\n"
        "예수님의 이름으로 기도드립니다. 아멘."
    ),
}


if __name__ == "__main__":
    template = build_template("QT_묵상_템플릿_v2.pptx")
    generate_qt_pptx(SAMPLE_DATA, template, "QT_출력_결과.pptx")
